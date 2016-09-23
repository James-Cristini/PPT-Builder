import os
import sys
from PIL import Image
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

"""
Initial thoughts:

Organize directories into names, countries, and file names

First need to get list of the names
-This can be done either through getting input via GUI or by reading a file and pulling out names via REGEX

Walk through the HW dictionary and paste images into either a single image or directy onto a powerpoint slide

Would it be worthwhile to rearrange files by name instead?

Create a new folder which contains all of the master images for each name/country

Need a GUI?

Potential Stretch:
Discuss data source format and if data for each name can be auto-populated into the deck
"""

def get_list_of_names(start_dir, country_list):
    names_from_files = []

    for x in country_list:
        country_path = start_dir + x
        file_list = [x for x in os.listdir(country_path)]
        for i in file_list:
            line = i.split("-")
            if i[-3:] == "jpg":
                line2 = line[0].split("_")
                if line2[1] not in names_from_files:
                    names_from_files.append(line2[1])

    return names_from_files


start_dir = os.getcwd() + "\HW\\"
country_list = [x for x in os.listdir(start_dir)]
names_list = get_list_of_names(start_dir, country_list)
HW = {}

for c in country_list:
    HW[c] = {}
    for n in names_list:
        HW[c][n] = []

# Creates a dictionary mapping out countries, names, and file names for each

for c in HW:
    #print c
    for n in HW[c]:
        #print n
        count = 1
        c_range = 5
        if c == "USA" or c == "EU":
            c_range = 10
        for x in range(0, c_range):
            if count == 10:
                file_name = c + "_" + n + "-10.jpg"
            else:
                file_name = c + "_" + n + "-0" + str(count) + ".jpg"
            HW[c][n].append(file_name)
            count += 1


prs_slides = {}

prs = Presentation('ppts/aw_template.pptx') # Opens AW master template
hw_slide_layout = prs.slide_layouts[5] # Slide used for handwriting samples


for name in names_list:
    for country in country_list:

        count = 0

        print name, country
        sld = "Slide %d" % (x)
        prs_slides[sld] = prs.slides.add_slide(hw_slide_layout)
        title = prs_slides[sld].placeholders[18] #Title text in Red
        subtitle = prs_slides[sld].placeholders[16] # Subtitle text in Grey

        title.text = "HANDWRITING SAMPLES"
        subtitle.text = name
        textbox = prs_slides[sld].shapes[2]
        sp = textbox.element
        sp.getparent().remove(sp)

        ### Add grey box

        left = Inches(1.25)
        top = Inches(1.5)
        ht = Inches(4.5)
        wd = Inches(7.5)
        prs_slides[sld].shapes.add_picture("imgs/grey_box.png", left, top, wd, ht)

        ### Add flag image

        flag_img = "imgs/flag_%s.png" % (country)
        left = Inches(8)
        top = Inches(0.7)
        ht = Inches(0.7)
        wd = Inches(0.7)
        prs_slides[sld].shapes.add_picture(flag_img, left, top, wd, ht)

        ### Add HW samples to each
        for x in range(len(HW[country][name])):
            img_file = "HW\\" + country + "\\" + HW[country][name][x]
            im = Image.open(img_file)
            ht = Inches(1.25)
            wd = Inches(2.0)
            #print img_file
            # im.show()
            # raw_input("Enter")

            if "01" in img_file or "06" in img_file or "02" in img_file:
                top = Inches(1.75)
            if "07" in img_file or "03" in img_file or "08" in img_file:
                top = Inches(3.125)
            if "04" in img_file or "09" in img_file or "05" in img_file or "10" in img_file:
                top = Inches(4.5)

            if "01" in img_file or "07" in img_file or "04" in img_file:
                left = Inches(1.5)
            if "06" in img_file or "03" in img_file or "09" in img_file:
                left = Inches(4.0)
            if "02" in img_file or "08" in img_file or "05" in img_file or "10" in img_file:
                left = Inches(6.5)

            prs_slides[sld].shapes.add_picture(img_file, left, top, wd, ht)


prs.save('hw_test.pptx')
print "DONE"


"""
    left = Inches(1)
    top = Inches(1.4)
    ht = Inches(4.5)
    wd = Inches(8)
    prs_slides[sld].shapes.add_picture("imgs/grey_box.png", left, top, wd, ht)

    ### Add flag image

    flag_img = "imgs/flag_%s.png" % (country)
    left = Inches(8)
    top = Inches(0.7)
    ht = Inches(0.7)
    wd = Inches(0.7)
    prs_slides[sld].shapes.add_picture(flag_img, left, top, wd, ht)

    ### Add HW samples to each
    for x in range(len(HW[country][name])):
        img_file = "HW\\" + country + "\\" + HW[country][name][x]
        im = Image.open(img_file)
        ht = Inches(1.47)
        wd = Inches(2.1)
        #print img_file
        # im.show()
        # raw_input("Enter")

        if "01" in img_file:
            #print "1", img_file
            left = Inches(1.5)
            top = Inches(1.5)

        elif "02" in img_file:
            #print "2", img_file
            left = Inches(6.5)
            top = Inches(1.5)

        elif "03" in img_file:
            #print "3", img_file
            left = Inches(4)
            top = Inches(2.75)

        elif "04" in img_file:
            #print "4", img_file
            left = Inches(1.5)
            top = Inches(4)

        elif "05" in img_file:
            #print "5", img_file
            left = Inches(6.5)
            top = Inches(4)

        elif "06" in img_file:
            left = Inches(4)
            top = Inches(1.5)

        elif "07" in img_file:
            left = Inches(1.5)
            top = Inches(2.75)

        elif "08" in img_file:
            left = Inches(6.5)
            top = Inches(2.75)

        elif "09" in img_file:
            left = Inches(4)
            top = Inches(4)

        #10 can go over 05's position
        elif "10" in img_file:
            left = Inches(6.5)
            top = Inches(4)


        prs_slides[sld].shapes.add_picture(img_file, left, top, wd, ht)

"""
