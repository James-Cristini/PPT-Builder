
"""
The ppt builder is intended for use by the Market Research department at Addison Whitney.

The program assists in putting together handwriting sample slides for brand name evaluation decks.
Takes in a folder of handwriting samples for potential brand names organized by country and builds a number of
slides in a matter of minutes

Please note that handwriting samples are not included for confidentiality purposes

This program is intended for use on Windows
"""

__author__ = "James Cristini"
__credits__ = "James Cristini"
__version__ = "1.3"
__maintainer__ = "James Cristini"
__email__ = "jacristi0428@gmail.com"

import os
import sys
import sip
from PyQt4 import QtGui, QtCore, uic
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

class MainWindow(QtGui.QMainWindow):

    INSTRUCTION_TEXT = """Copy the entire HW folder (and make sure the folder is called HW)
to the ppt_builder folder before pressing the "Build Slides" button!

Be sure your folders are structured as the example below:

ppt_builer (parent folder containing the data folder and ppt_builder.exe)
 - HW (folder containing country folders - remember this folder must be named HW)
   - AUS (country folder containing handwriting sample image files)
   - USA (be sure each of these folders are named for the 3-letter country code)
   - CAN
   - JPN


The value in the number box in the bottom left should match the slide number of the
desired slide to use for samples in the ppt template's slide master view
This should not change unless the template iteself has changed
"""

    def __init__(self):
        super(MainWindow, self).__init__()
        """ Main window for the powerpoint builder """

        #Loads the UI file; kept the UI file separate since it can be easily stored in the data folder (easier to update if it remains separate)
        self.ui = uic.loadUi('data/UI.ui')
        self.ui.setWindowIcon(QtGui.QIcon("data/pencil.ico"))

        self.ui.progress_bar.setValue(0)

        self.ui.text_output.setPlainText(self.INSTRUCTION_TEXT)

        self.ui.start_btn.clicked.connect(self.build_deck)

        self.remove_qt_conf()

        self.ui.show()

    def remove_qt_conf(self):
        """ PyQt generates an uneeded qt.conf file that can be reomoved if it is there """
        try:
            os.remove('qt.conf')
        except WindowsError:
            # Just pass in the event that the file is not there
            pass

    # Parses out brand names from image file names
    def get_list_of_names(self, start_dir, country_list):
        """ Parses out brand names from image files found through start_dir + country name folders """
        names_from_files = []

        # This is specific to the current file name conventions
        for x in country_list:
            country_path = start_dir + x
            file_list = [x for x in os.listdir(country_path)]
            for i in file_list:
                line = i.split("-")
                if i[-3:] == "jpg":
                    line2 = line[0].split("_")
                    if line2[1] not in names_from_files:
                        names_from_files.append(line2[1])

        return names_from_files

    def build_deck(self) :
        """ Build the powerpoint deck """
        try: #wrapped in a big try/except to account for and output any unexpected errors
            start_dir = os.getcwd() + "\HW\\"
            country_list = [x for x in os.listdir(start_dir)]
            names_list = self.get_list_of_names(start_dir, country_list)

            HW = {}

            total = len(country_list) * len(names_list)
            self.ui.progress_bar.setMaximum(total)

            # Initialize the dictionary that will contain {country_code: {name_candidate : file_name}}
            for c in country_list:
                HW[c] = {}
                for n in names_list:
                    HW[c][n] = []

            # Creates a dictionary mapping out countries, names, and file names for each
            # Sets prority on certain countries in terms of placement in the PPT
            # This is again specific to current file name conventions
            for c in HW:
                for n in HW[c]:
                    count = 1
                    c_range = 5
                    if c == "USA" or c == "EU":
                        c_range = 10
                    for x in range(0, c_range):
                        if count == 10:
                            file_name = c + "_" + n + "-10.jpg"
                        else:
                            file_name = c + "_" + n + "-0" + str(count) + ".jpg"
                        HW[c][n].append(file_name)
                        count += 1

            # Initialize a dictionary that will store the new PPT deck slides
            prs_slides = {}

            ### ========================
            ### Adjust the below file name to reflect the template name and path (preferably keep the same path, e.g. data/{template_name})
            ### ========================
            prs = Presentation('data/aw_template.pptx') # Opens AW master template

            ### ========================
            ### prs.slide_layouts[] will need to be adjusted in the event of a template change, it is the 15th item in the list of slides
            ### when counting, skip the first (larger) slide and start the following slide at 0
            ### Or count them all starting at 1 and subtract 2 ;)
            ### ========================
            slide_num = self.ui.slide_to_use.value()
            hw_slide_layout = prs.slide_layouts[slide_num]
            placeholders = hw_slide_layout.placeholders

            ### ========================
            ### The below code sets the placeholders that are used for the Title and Subtitle on each slide
            ### placeholder_idx_list below print the idx number of each placeholder on the given slide
            ### This may or may not always work - I assume if the slide has only 2 placeholders, it should always work just fine
            ### ========================
            placeholder_idx_list = [i.placeholder_format.idx for i in placeholders]
            try:
                slide_title = placeholder_idx_list[0]
                slide_subtitle = placeholder_idx_list[1]
            except IndexError:
                self.ui.text_output.setPlainText("This slide does not have a title and/or subtitle to draw from")
                return

            ### ========================
            ### Sets an error image in the case that an image cannot be found
            ### The image itself can be replaced as long as the path/name remains the same
            ### ========================
            ERROR_IMG = "data/error_img.png"

            # Progress bar starting value
            bar_val = 0

            #Text to output in the UI while building deck, starts blank but will be added to as events are processed
            out_text = ""

            # Build slides looking for at names, then building slides for each country associated with the name
            for name in names_list:
                for country in country_list:

                    # Updated UI output text changes as slides are built, update progress bar accordingly as well
                    out_text = "{} {} slide added\n".format(name, country)
                    self.ui.text_output.setPlainText(out_text)
                    self.ui.progress_bar.setValue(bar_val)

                    # New slide is created and added to the pr_slides dictionary / PPT deck
                    sld = "Slide {}".format(x)
                    prs_slides[sld] = prs.slides.add_slide(hw_slide_layout)

                    title = prs_slides[sld].placeholders[slide_title] #Title text in Red
                    subtitle = prs_slides[sld].placeholders[slide_subtitle] # Subtitle text in Grey

                    title.text = "Handwriting samples"
                    subtitle.text = name

                    # Add flag image
                    flag_img = "data/flag_{}.png".format(country)
                    left = Inches(9.0) # x inches from the left
                    top = Inches(0.5) # x inches from the top
                    ht = Inches(0.5) # obj height
                    wd = Inches(0.5) # obj width
                    try:
                        prs_slides[sld].shapes.add_picture(flag_img, left, top, wd, ht)
                    except: # Add error image if the flag is not found
                        prs_slides[sld].shapes.add_picture(ERROR_IMG, left, top, wd, ht)

                    # Add HW sample iamges to each
                    for x in range(len(HW[country][name])):
                        img_file = "HW\\" + country + "\\" + HW[country][name][x]
                        ht = Inches(1.25)
                        wd = Inches(2.0)

                        # Position based on the number found in the image file name
                        if "01" in img_file or "06" in img_file or "02" in img_file:
                            top = Inches(1.3)
                        if "07" in img_file or "03" in img_file or "08" in img_file:
                            top = Inches(2.55)
                        if "04" in img_file or "09" in img_file or "05" in img_file or "10" in img_file:
                            top = Inches(3.8)

                        if "01" in img_file or "07" in img_file or "04" in img_file:
                            left = Inches(1.5)
                        if "06" in img_file or "03" in img_file or "09" in img_file:
                            left = Inches(4.0)
                        if "02" in img_file or "08" in img_file or "05" in img_file or "10" in img_file:
                            left = Inches(6.5)
                        try:
                            prs_slides[sld].shapes.add_picture(img_file, left, top, wd, ht)
                        except: # Add error image if the expected sample is not found
                            prs_slides[sld].shapes.add_picture(ERROR_IMG, left, top, wd, ht)

                    # Update progress bar value and process UI events so that the progress bar updates are seen per iter
                    bar_val += 1
                    QtGui.QApplication.processEvents()

            try: # Save the file and open. If the old file is already open, show a message box detailing the problem
                prs.save('hw_slides.pptx')
                self.open_deck('hw_slides.pptx')
            except IOError:
                QtGui.QMessageBox.warning(self, "File already open", "The hw_slides file is currently open; close this file first and try again")

            # Final updates to progress bar and output text
            self.ui.progress_bar.setValue(total)
            self.ui.text_output.setPlainText("All slides have been built!")

        # Wrapped the entire function in a try/except to easily show any unexpected errors that might come up when running easier debugging ;)
        except Exception as e:
            self.ui.text_output.setPlainText(str(e))

    def open_deck(self, file_name):
        """ opens the ppt deck after completion """
        # Prompt the user to open the ppt, if yes, open - otherwise just state that it has been saved
        choice = QtGui.QMessageBox.question(self, "Open", "The powerpoint deck has been built, would you like to open now?", QtGui.QMessageBox.Yes | QtGui.QMessageBox.No)

        if choice ==  QtGui.QMessageBox.Yes:
            try:
                os.system("start " + str(file_name))
            except IOError:
                QtGui.QMessageBox.warning(self, "File already open", "The hw_slides file is currently open; close this file first and try again") % ( str(self.filter_file_name))
        else:
            QtGui.QMessageBox.warning(self, "Saved", "The file has been saved as hw_slides")
            self.ui.text_output.setPlainText("The completed deck has been saved as hw_slides")

if __name__ == '__main__':
    app = QtGui.QApplication(sys.argv)
    app.setStyle("Plastique")
    app.setWindowIcon(QtGui.QIcon("book.ico"))
    window = MainWindow()

    sip.setdestroyonexit(False)
    sys.exit(app.exec_())
