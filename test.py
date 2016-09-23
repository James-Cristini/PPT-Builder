import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

#TODO be sure all flags are gotten with proper abreviations
#TODO get more precise shape placements

#TODO separate slides for USA
#TODO separate slides for EU5? would it be better to keep them separate with 5 per?
#TODO DELETE empty palceholders...
#TODO Build basic UI WITH progress bar



"""
def analyze_ppt(input, output):

    prs = Presentation(input)
    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = 'Title for Layout {}'.format(index)
        except AttributeError:
            print("No Title for Layout {}".format(index))
        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if 'Title' not in shape.text:
                        shape.text = 'Placeholder index:{} type:{}'.format(phf.idx, shape.name)
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                print('{} {}'.format(phf.idx, shape.name))
    prs.save(output)

analyze_ppt("aw_template.pptx", "test.pptx")
"""

prs_slides = {}

prs = Presentation('ppts/aw_template.pptx') # Opens AW master template
hw_slide_layout = prs.slide_layouts[5] # Slide used for handwriting samples
# slide = prs.slides.add_slide(hw_slide_layout) # Adds a new blank slide to be modified
#
# title = slide.placeholders[18] #Title text in Red
# subtitle = slide.placeholders[16] # Subtitle text in Grey


for x in range(0,5):
    sld = "Slide %d" % (x)
    prs_slides[sld] = prs.slides.add_slide(hw_slide_layout)
    title = prs_slides[sld].placeholders[18] #Title text in Red
    subtitle = prs_slides[sld].placeholders[16] # Subtitle text in Grey

    #body = prs_slides[sld].placeholders[17] # NEED TO DELETE THIS ONE

    textbox = prs_slides[sld].shapes[2]
    sp = textbox.element
    sp.getparent().remove(sp)

    title.text = "Slide %d" % (x)
    subtitle.text = "This is slide %d" % (x)

for x in prs_slides:
    print x, prs_slides[x]

# title.text = "Handwriting Samples"
# subtitle.text = "Adexli"
#
# left = Inches(1)
# top = Inches(1.4)
# ht = Inches(4.5)
# wd = Inches(8)
# slide.shapes.add_picture("imgs/grey_box.png", left, top, wd, ht)
#
# ###
#
# left = Inches(8)
# top = Inches(.8)
# ht = Inches(0.47)
# wd = Inches(0.71)
# slide.shapes.add_picture("imgs/flag_AUS.png", left, top, wd, ht)
#
# ###
#
# left = Inches(1.5)
# top = Inches(1.5)
# ht = Inches(1.47)
# wd = Inches(2.14)
# slide.shapes.add_picture("HW/AUS/AUS_Adexli-01.jpg", left, top, wd, ht)
#
# left = Inches(6.5)
# top = Inches(1.5)
# ht = Inches(1.47)
# wd = Inches(2.14)
# slide.shapes.add_picture("HW/AUS/AUS_Adexli-02.jpg", left, top, wd, ht)
#
# left = Inches(4)
# top = Inches(2.75)
# ht = Inches(1.47)
# wd = Inches(2.14)
# slide.shapes.add_picture("HW/AUS/AUS_Adexli-03.jpg", left, top, wd, ht)
#
# left = Inches(1.5)
# top = Inches(4)
# ht = Inches(1.47)
# wd = Inches(2.14)
# slide.shapes.add_picture("HW/AUS/AUS_Adexli-04.jpg", left, top, wd, ht)
#
# left = Inches(6.5)
# top = Inches(4)
# ht = Inches(1.47)
# wd = Inches(2.14)
# slide.shapes.add_picture("HW/AUS/AUS_Adexli-05.jpg", left, top, wd, ht)



prs.save('test.pptx')
